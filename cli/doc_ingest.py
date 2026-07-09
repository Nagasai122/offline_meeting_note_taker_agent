"""
Offline document context ingestion: PDF/PPTX/DOCX/XLSX/TXT/images ->
extracted text -> LLM-summarised context (bounded to ~1000 tokens per
attachment, ~1500 tokens combined) -> `.doc_context.txt`.

All extraction libraries (pdfplumber, python-pptx, python-docx, openpyxl,
pytesseract) are pure Python (pytesseract shells out to a local Tesseract
binary, never a network call) and read local files only -- no network calls,
consistent with the zero-egress guarantee the rest of this project is built
around. The 1000-token cap on each attachment's summary is a deliberate
trade-off (architecture_v2.md §7.2): complete coverage of e.g. a 200-slide
deck is worth less than keeping the extraction context window clear for the
actual transcript.

P3 (live-recording "Add Context"): context can now arrive at any point during
a session, not just pre-meeting, so `add_document_context`/
`add_pasted_text_context` APPEND a labelled section per source rather than
overwriting the file -- a second attachment must not silently discard the
first. Every write goes through a FileLock (the same global lock every other
writer in this project uses, per concurrency/lock.py's convention) since two
attachments added in quick succession during a live recording could otherwise
race on a read-modify-write of the same sidecar file.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Callable

from concurrency.atomic import atomic_write_text
from concurrency.lock import FileLock

DOC_SUMMARY_SYSTEM_PROMPT = (
    "Summarise the following section of a meeting document in 3-5 bullet points. "
    "Focus on key arguments, data, decisions, and terminology."
)

_MAX_OUTPUT_CHARS_PER_TOKEN = 4  # rough token->char approximation for the truncation cap

# P3: attachments shorter than this aren't worth an LLM summarisation call --
# the round trip costs more than the (negligible) compression it would buy.
# Same "don't pay for work below a meaningful threshold" philosophy as
# mcp_server/tools/extraction.py's MIN_TRANSCRIPT_WORDS.
_MIN_WORDS_FOR_SUMMARY = 60

# P3.8 (cap-and-compress): once the combined accumulated doc_context.txt
# exceeds this many characters (~1500 tokens), the whole file is re-summarised
# into one consolidated block instead of growing unbounded -- N attachments
# must never claim more than roughly one attachment's worth of the
# extraction prompt's tight token budget.
_COMBINED_CONTEXT_CHAR_CAP = 1500 * _MAX_OUTPUT_CHARS_PER_TOKEN


def extract_text_from_pdf(path: Path) -> str:
    import pdfplumber

    try:
        with pdfplumber.open(path) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
    except Exception as exc:  # pdfplumber raises varied exception types for encrypted PDFs
        if "password" in str(exc).lower() or "encrypt" in str(exc).lower():
            raise ValueError("PDF is encrypted") from exc
        raise
    return "\n\n".join(pages)


def extract_text_from_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    parts.append(text)
        slides_text.append("\n".join(parts))
    return "\n\n---\n\n".join(slides_text)


def extract_text_from_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    return "\n".join(p.text for p in document.paragraphs)


def extract_text_from_xlsx(path: Path) -> str:
    """Row-by-row, sheet-by-sheet text dump -- good enough for the same
    "context for the LLM", not "faithful reproduction" bar every other
    extractor here holds itself to. Cells are joined with tabs (visually
    approximates the grid without pulling in a table-formatting dependency),
    blank rows are skipped so a sparse sheet doesn't pad the raw text with
    hundreds of empty lines before it ever reaches the word-count-based
    chunker."""
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets_text = []
    for sheet in workbook.worksheets:
        rows_text = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append("\t".join(cells))
        if rows_text:
            sheets_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text))
    return "\n\n---\n\n".join(sheets_text)


_ocr_available_cache: bool | None = None


def is_ocr_available() -> bool:
    """Health-check for the local Tesseract OCR binary -- mirrors
    llm/http_probe.py's "probe before you need it" pattern so a missing
    install surfaces as a clear status flag (see GET /api/server/status)
    rather than only as a cryptic error the first time someone happens to
    upload an image. Cached after the first check: installation status
    cannot change during one run, and this is polled from the dashboard's
    System tab every few seconds -- spawning `tesseract --version` on every
    poll would be wasted work for a fact that only needs checking once.
    """
    global _ocr_available_cache
    if _ocr_available_cache is None:
        try:
            import pytesseract

            pytesseract.get_tesseract_version()
            _ocr_available_cache = True
        except Exception:
            _ocr_available_cache = False
    return _ocr_available_cache


def extract_text_from_image(path: Path) -> str:
    """Local OCR via pytesseract (a thin wrapper around the separate
    Tesseract OCR binary -- not installable via pip, see docs/setup-guide.md).
    Fully local/offline, consistent with the zero-egress posture: no cloud
    OCR API is ever involved.

    Raises:
        RuntimeError: if the Tesseract binary isn't installed/on PATH --
            surfaced as a clear, actionable error rather than pytesseract's
            own cryptic TesseractNotFoundError, since this is the one
            extractor in this module with an external-binary dependency the
            other pure-Python extractors don't have.
    """
    import pytesseract
    from PIL import Image

    try:
        with Image.open(path) as img:
            text = pytesseract.image_to_string(img)
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError(
            "OCR unavailable: the Tesseract OCR binary is not installed or not on PATH. "
            "See docs/setup-guide.md for the Windows installer link, or set "
            "pytesseract.pytesseract.tesseract_cmd to its install location."
        ) from exc
    return text


def extract_text(path: Path) -> str:
    """Dispatch to the right extractor based on file suffix.

    Args:
        path: Path to the uploaded document.

    Returns:
        Raw extracted text.

    Raises:
        NotImplementedError: for legacy .ppt/.doc binary formats.
        ValueError: for any other unsupported extension.
        RuntimeError: for an image file when the Tesseract OCR binary isn't
            installed (see extract_text_from_image).
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    if suffix == ".pptx":
        return extract_text_from_pptx(path)
    if suffix == ".ppt":
        raise NotImplementedError("Legacy .ppt is not supported; please save as .pptx.")
    if suffix == ".docx":
        return extract_text_from_docx(path)
    if suffix == ".doc":
        raise NotImplementedError("Legacy .doc is not supported; please save as .docx.")
    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".xlsx":
        return extract_text_from_xlsx(path)
    if suffix in (".png", ".jpg", ".jpeg"):
        return extract_text_from_image(path)
    raise ValueError(f"Unsupported document extension: {suffix!r}")


def _chunk_text(text: str, chunk_tokens: int = 2000) -> list[str]:
    """Split text into ~chunk_tokens-sized pieces by word count (reuses the
    same words*1.35 approximation as transcribe/chunker.py, applied to plain
    prose rather than transcript segments)."""
    words = text.split()
    if not words:
        return []
    words_per_chunk = max(1, int(chunk_tokens / 1.35))
    return [
        " ".join(words[i : i + words_per_chunk])
        for i in range(0, len(words), words_per_chunk)
    ]


def summarise_doc_context(
    raw_text: str,
    llm_call: Callable[[str, str], str],
    max_output_tokens: int = 1000,
) -> str:
    """Summarise `raw_text` in ~2000-token chunks, then cap total output.

    Args:
        raw_text: Extracted document text.
        llm_call: Callable(system_prompt, user_text) -> response text.
        max_output_tokens: Approximate token budget for the final summary.

    Returns:
        Concatenated bullet-point summary, truncated (on a bullet boundary)
        to roughly `max_output_tokens`.
    """
    chunks = _chunk_text(raw_text, chunk_tokens=2000)
    if not chunks:
        return ""

    summaries = [llm_call(DOC_SUMMARY_SYSTEM_PROMPT, chunk) for chunk in chunks]
    combined = "\n".join(summaries)

    max_chars = max_output_tokens * _MAX_OUTPUT_CHARS_PER_TOKEN
    if len(combined) <= max_chars:
        return combined

    # Truncate on a bullet-point boundary rather than mid-line.
    lines = combined.splitlines()
    kept: list[str] = []
    total = 0
    for line in lines:
        if total + len(line) + 1 > max_chars:
            break
        kept.append(line)
        total += len(line) + 1
    return "\n".join(kept)


def _is_worth_summarising(text: str) -> bool:
    return len(text.split()) >= _MIN_WORDS_FOR_SUMMARY


def _append_labelled_context(
    session_id: str,
    meetings_dir: Path,
    source_label: str,
    body: str,
    llm_call: Callable[[str, str], str],
    lock_path: Path | str,
    lock_timeout: float,
) -> Path:
    """Append one labelled section to the session's accumulating
    `<session_id>.doc_context.txt` sidecar under a FileLock.

    A blank/whitespace-only body (e.g. an empty document, or OCR that found
    no text) writes nothing and leaves the file untouched -- there's no value
    in a labelled section with no content, and creating one would permanently
    clutter the extraction context with noise.
    """
    meetings_dir = Path(meetings_dir)
    meetings_dir.mkdir(parents=True, exist_ok=True)
    output_path = meetings_dir / f"{session_id}.doc_context.txt"

    body = body.strip()
    if not body:
        return output_path

    timestamp = datetime.now().strftime("%H:%M")
    entry = f"--- {source_label} ({timestamp}) ---\n{body}\n"

    with FileLock(lock_path, timeout_seconds=lock_timeout):
        existing = output_path.read_text(encoding="utf-8") if output_path.exists() else ""
        combined = f"{existing}\n{entry}" if existing else entry
        if len(combined) > _COMBINED_CONTEXT_CHAR_CAP:
            combined = summarise_doc_context(combined, llm_call, max_output_tokens=1500)
        atomic_write_text(output_path, combined)
    return output_path


def add_document_context(
    path: Path,
    session_id: str,
    meetings_dir: Path,
    llm_call: Callable[[str, str], str],
    lock_path: Path | str,
    lock_timeout: float,
    source_label: str | None = None,
) -> Path:
    """Full pipeline for one uploaded document: extract -> (maybe) summarise
    -> append (labelled by filename) to `<session_id>.doc_context.txt`.

    Renamed from the old `ingest_document`, which overwrote the file on every
    call -- fine when context could only be attached once, pre-meeting, but
    P3 allows adding context throughout a live recording, where a second
    attachment must accumulate alongside the first rather than silently
    discard it.

    Args:
        path: Path to the uploaded document on disk (PDF/PPTX/DOCX/XLSX/TXT/
            PNG/JPG) -- often a tmp-staged copy under a caller-chosen name,
            not necessarily the name the user sees.
        session_id: Session this context belongs to.
        meetings_dir: Directory the `.doc_context.txt` artefact is written into.
        llm_call: Callable(system_prompt, user_text) -> response text.
        lock_path: Concurrency lock path (same global lock every other writer
            in this project uses -- see concurrency/lock.py's module docstring).
        lock_timeout: Lock acquisition timeout in seconds.
        source_label: The filename to show in the labelled section (e.g. the
            original upload filename). Defaults to `path.name` when omitted,
            but callers that stage the upload under a different on-disk name
            (cli/web.py's upload endpoint prefixes it with `<session_id>_doc_`
            to avoid collisions) must pass the real filename explicitly here
            -- otherwise the label leaks the tmp-staging naming scheme into
            the extraction prompt instead of the name the user actually gave
            the file.

    Returns:
        Path to the (possibly just-appended-to) `.doc_context.txt` file.
    """
    raw_text = extract_text(path)
    body = summarise_doc_context(raw_text, llm_call) if _is_worth_summarising(raw_text) else raw_text
    return _append_labelled_context(
        session_id, meetings_dir, f"Document: {source_label or path.name}", body,
        llm_call, lock_path, lock_timeout,
    )


def add_pasted_text_context(
    text: str,
    session_id: str,
    meetings_dir: Path,
    llm_call: Callable[[str, str], str],
    lock_path: Path | str,
    lock_timeout: float,
) -> Path:
    """Same accumulation pipeline as add_document_context, for text pasted
    directly (e.g. a Teams chat snippet) rather than uploaded as a file."""
    body = summarise_doc_context(text, llm_call) if _is_worth_summarising(text) else text
    return _append_labelled_context(
        session_id, meetings_dir, "Pasted text", body, llm_call, lock_path, lock_timeout,
    )
